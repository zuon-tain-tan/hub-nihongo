import os
from pptx import Presentation
from docx import Document
from docx.shared import Cm, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
from docx.oxml.ns import qn
import io

def set_font_style(run, font_name='Times New Roman', east_asia_font='Yu Gothic', size=12, color=RGBColor(0, 0, 0), bold=False):
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    r = run._element
    rPr = r.get_or_add_rPr()
    rFonts = rPr.get_or_add_rFonts()
    rFonts.set(qn('w:eastAsia'), east_asia_font)

def get_image_from_shape(shape):
    """Lấy dữ liệu ảnh từ nhiều loại shape khác nhau"""
    if shape.shape_type == 13: 
        return io.BytesIO(shape.image.blob)
    if shape.shape_type == 14:
        try:
            if shape.image:
                return io.BytesIO(shape.image.blob)
        except:
            pass
    try:
        if shape.fill.type == 6: # msoFillPicture
            return io.BytesIO(shape.fill.picture.image.blob)
    except:
        pass
    return None

def create_review_word():
    doc = Document()
    
    # Cấu hình trang A4
    section = doc.sections[0]
    section.page_height, section.page_width = Cm(29.7), Cm(21.0)
    section.left_margin = section.right_margin = section.top_margin = section.bottom_margin = Cm(2.0)

    heading = doc.add_heading('ÔN TẬP TỪ VỰNG BÀI 19 ĐẾN 25', level=1)
    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in heading.runs:
        set_font_style(run, size=16, bold=True)

    ppt_path = r'ppt/bai19den25.pptx'
    if not os.path.exists(ppt_path):
        print(f"Không tìm thấy file {ppt_path}")
        return

    prs = Presentation(ppt_path)
    table = doc.add_table(rows=0, cols=2)
    table.style = 'Table Grid'
    table.autofit = False
    
    for slide in prs.slides:
        image_stream = None
        slide_text = ""
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text: slide_text += text + " "
            
            if image_stream is None:
                image_stream = get_image_from_shape(shape)

        row_cells = table.add_row().cells
        
        # --- Căn giữa theo chiều dọc cho cả 2 ô (Align Center) ---
        row_cells[0].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        row_cells[1].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER

        # Ô 1: Chèn ảnh (Chiều cao 4cm)
        if image_stream:
            p_img = row_cells[0].paragraphs[0]
            p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Cấu hình khoảng cách cho ô ảnh (Line 1, Before/After 6pt)
            p_img.paragraph_format.space_before = Pt(6)
            p_img.paragraph_format.space_after = Pt(6)
            p_img.paragraph_format.line_spacing = 1.0

            run_img = p_img.add_run()
            try:
                run_img.add_picture(image_stream, height=Cm(4)) # Đã đổi sang xét chiều cao
            except:
                p_img.add_run("[Lỗi định dạng ảnh]")
        
        # Ô 2: Chèn text tiếng Nhật (Căn trái)
        p_txt = row_cells[1].paragraphs[0]
        p_txt.alignment = WD_ALIGN_PARAGRAPH.LEFT # Căn lề trái
        
        # Cấu hình khoảng cách cho ô chữ (Line 1, Before/After 6pt)
        p_txt.paragraph_format.space_before = Pt(6)
        p_txt.paragraph_format.space_after = Pt(6)
        p_txt.paragraph_format.line_spacing = 1.0

        run_txt = p_txt.add_run(slide_text.strip())
        set_font_style(run_txt, font_name='Times New Roman', east_asia_font='Yu Gothic', size=12)

    doc.save('On_tap_tu_vung_19_den_25.docx')
    print("Đã tạo file thành công: Ảnh cao 4cm, căn Align Center Left và dãn dòng 6pt!")

if __name__ == "__main__":
    create_review_word()